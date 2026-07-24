"""
统一的文件解析工具
支持 PDF, PPTX, DOCX, MD, TXT 等格式
"""
import os
from typing import Optional


def parse_file(file_path: str) -> str:
    """
    根据文件扩展名自动选择解析器，返回纯文本内容
    """
    ext = os.path.splitext(file_path)[1].lower()

    if ext == '.pdf':
        return _parse_pdf(file_path)
    elif ext == '.pptx':
        return _parse_pptx(file_path)
    elif ext == '.ppt':
        return _parse_ppt(file_path)
    elif ext == '.docx':
        return _parse_docx(file_path)
    elif ext == '.doc':
        return _parse_doc(file_path)
    elif ext in ('.txt', '.md', '.cpp', '.h', '.hpp', '.py', '.json', '.xml', '.yaml', '.yml'):
        return _parse_text(file_path)
    else:
        raise ValueError(f"不支持的文件类型: {ext}")


def _parse_pdf(file_path: str) -> str:
    try:
        from pypdf import PdfReader
        reader = PdfReader(file_path)
        texts = []
        for page in reader.pages:
            text = page.extract_text()
            if text:
                texts.append(text)
        return "\n\n".join(texts)
    except ImportError:
        raise ImportError("请安装 pypdf: pip install pypdf")


def _parse_pptx(file_path: str) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        texts = []
        for slide in prs.slides:
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            slide_texts.append(para.text.strip())
            if slide_texts:
                texts.append("\n".join(slide_texts))
        return "\n\n---\n\n".join(texts)
    except ImportError:
        raise ImportError("请安装 python-pptx: pip install python-pptx")


def _parse_ppt(file_path: str) -> str:
    """旧 .ppt 格式（OLE2），python-pptx 不支持。尝试读取嵌入文本"""
    try:
        # 尝试用 olefile 读取（如果安装了的话）
        import olefile
        ole = olefile.OleFileIO(file_path)
        # 尝试从 PowerPoint Document 流中提取文本
        if ole.exists('PowerPoint Document'):
            data = ole.openstream('PowerPoint Document').read()
            # 简单提取可读 ASCII/UTF-8 文本
            text = data.decode('utf-8', errors='ignore')
            # 过滤不可打印字符
            import re
            clean = re.sub(r'[^\x20-\x7e一-鿿　-〿＀-￯\n\r\t]', '', text)
            if len(clean) > 100:
                return clean[:10000]
        ole.close()
    except ImportError:
        pass
    except Exception:
        pass
    raise ValueError("旧版 .ppt 格式暂不支持文本解析，请转换为 .pptx 后上传")


def _parse_doc(file_path: str) -> str:
    """旧 .doc 格式，python-docx 不支持"""
    raise ValueError("旧版 .doc 格式暂不支持文本解析，请转换为 .docx 后上传")


def _parse_docx(file_path: str) -> str:
    try:
        from docx import Document
        doc = Document(file_path)
        texts = []
        for para in doc.paragraphs:
            if para.text.strip():
                texts.append(para.text.strip())
        return "\n\n".join(texts)
    except ImportError:
        raise ImportError("请安装 python-docx: pip install python-docx")


def _parse_text(file_path: str) -> str:
    encodings = ['utf-8', 'gbk', 'gb2312', 'latin-1']
    for enc in encodings:
        try:
            with open(file_path, 'r', encoding=enc) as f:
                return f.read()
        except UnicodeDecodeError:
            continue
    raise ValueError(f"无法解码文件: {file_path}")
